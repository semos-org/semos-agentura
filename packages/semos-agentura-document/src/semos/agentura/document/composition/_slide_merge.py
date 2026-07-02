"""Merge cherry-picked slides from multiple PPTX sources.

Dual backend: COM (PowerPoint, preserves animations/media) primary,
python-pptx (portable, pure Python) fallback.
"""

from __future__ import annotations

import copy as copy_module
import logging
import shutil
from pathlib import Path

from pydantic import BaseModel

from ..exceptions import SlideMergeError
from ..models import MergeResult
from ._backends import select_merge_backend

logger = logging.getLogger(__name__)


# ------------------------------------------------------------------
# Config model
# ------------------------------------------------------------------


class SlideRef(BaseModel):
    """Reference to a single slide in a source PPTX."""

    source: str  # key into MergeConfig.sources
    index: int  # 0-based slide index


class MergeConfig(BaseModel):
    """Configuration for a slide merge operation."""

    base: str  # path to base PPTX (theme)
    sources: dict[str, str]  # name -> path
    slides: list[SlideRef]
    output: str = "merged.pptx"
    notes: dict[int, str] = {}  # slide position -> note text


def load_merge_config(config_path: Path | str) -> MergeConfig:
    """Load a MergeConfig from a YAML file.

    Resolves relative paths against the config file's directory.
    """
    try:
        import yaml
    except ImportError as exc:
        raise SlideMergeError("PyYAML required for YAML config: pip install PyYAML") from exc

    config_path = Path(config_path).resolve()
    with open(config_path, encoding="utf-8") as f:
        raw = yaml.safe_load(f)

    base_dir = config_path.parent

    def _resolve(p: str) -> str:
        pp = Path(p)
        if not pp.is_absolute():
            pp = base_dir / pp
        return str(pp.resolve())

    return MergeConfig(
        base=_resolve(raw["base"]),
        sources={k: _resolve(v) for k, v in raw["sources"].items()},
        slides=[SlideRef(**s) for s in raw["slides"]],
        output=raw.get("output", "merged.pptx"),
        notes={int(k): v for k, v in raw.get("notes", {}).items()},
    )


def parse_source_args(
    args: list[str],
    *,
    base: str | None = None,
    output: str = "merged.pptx",
) -> MergeConfig:
    """Parse CLI source arguments into a MergeConfig.

    Args format: "path.pptx:0-5", "path.pptx:3,7,12",
    "path.pptx" (all slides).

    Examples:
        parse_source_args(["a.pptx:0-5", "b.pptx:3,7"])
        parse_source_args(["a.pptx:0-5"], base="template.pptx")
    """
    sources: dict[str, str] = {}
    slides: list[SlideRef] = []

    for arg in args:
        if ":" in arg:
            path, range_str = arg.rsplit(":", 1)
            # Check it's a range, not a Windows drive letter
            if len(range_str) > 1 or not range_str.isalpha():
                indices = _parse_range(range_str)
            else:
                # Was just "C:" drive prefix, treat as no range
                path = arg
                indices = None
        else:
            path = arg
            indices = None

        path = str(Path(path).resolve())
        name = Path(path).stem
        # Deduplicate source names
        if name in sources and sources[name] != path:
            i = 2
            while f"{name}_{i}" in sources:
                i += 1
            name = f"{name}_{i}"
        sources[name] = path

        if indices is None:
            # Will need to discover slide count at merge time
            slides.append(SlideRef(source=name, index=-1))
        else:
            for idx in indices:
                slides.append(SlideRef(source=name, index=idx))

    base_path = base if base else next(iter(sources.values()))
    return MergeConfig(
        base=str(Path(base_path).resolve()),
        sources=sources,
        slides=slides,
        output=output,
    )


def _parse_range(s: str) -> list[int]:
    """Parse '0-5' or '3,7,12' into a list of ints."""
    indices: list[int] = []
    for part in s.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            indices.extend(range(int(start), int(end) + 1))
        else:
            indices.append(int(part))
    return indices


# ------------------------------------------------------------------
# COM backend
# ------------------------------------------------------------------


class _COMBackend:
    """PowerPoint COM automation. Full fidelity."""

    def __init__(self) -> None:
        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        self.ppt = win32com.client.Dispatch("PowerPoint.Application")

    def merge(self, config: MergeConfig, output_path: Path) -> int:
        output_abs = str(output_path.resolve())
        shutil.copy2(config.base, output_abs)

        dst = self.ppt.Presentations.Open(output_abs, WithWindow=False)
        while dst.Slides.Count > 0:
            dst.Slides(1).Delete()

        count = 0
        for i, spec in enumerate(config.slides):
            src_path = config.sources[spec.source]
            src_idx = spec.index
            dst.Slides.InsertFromFile(
                src_path,
                dst.Slides.Count,
                src_idx + 1,
                src_idx + 1,
            )
            count += 1
            logger.info(
                "[%2d] %s[%d] -> output[%d]",
                i,
                spec.source,
                src_idx,
                i,
            )

        for pos, note_text in config.notes.items():
            if pos < dst.Slides.Count:
                try:
                    s = dst.Slides(pos + 1)
                    np = s.NotesPage.Shapes.Placeholders(2)
                    np.TextFrame.TextRange.Text = note_text
                except Exception:
                    pass

        dst.Save()
        dst.Close()
        return count

    def close(self) -> None:
        try:
            self.ppt.Quit()
        except Exception:
            pass


# ------------------------------------------------------------------
# python-pptx backend
# ------------------------------------------------------------------


class _PptxBackend:
    """Pure Python fallback via python-pptx XML clone."""

    def merge(self, config: MergeConfig, output_path: Path) -> int:
        from pptx import Presentation

        dst = Presentation(config.base)
        _delete_all_slides(dst)

        src_cache: dict[str, object] = {}
        count = 0

        for i, spec in enumerate(config.slides):
            src_path = config.sources[spec.source]
            if src_path not in src_cache:
                src_cache[src_path] = Presentation(src_path)
            src_prs = src_cache[src_path]

            if spec.index >= len(src_prs.slides):
                logger.warning(
                    "%s[%d] out of range (%d slides)",
                    spec.source,
                    spec.index,
                    len(src_prs.slides),
                )
                continue

            _clone_slide(src_prs.slides[spec.index], src_prs, dst)
            count += 1
            logger.info(
                "[%2d] %s[%d] -> output[%d]",
                i,
                spec.source,
                spec.index,
                i,
            )

        for pos, note_text in config.notes.items():
            if pos < len(dst.slides):
                ns = dst.slides[pos].notes_slide
                ns.notes_text_frame.text = note_text

        dst.save(str(output_path))
        return count

    def close(self) -> None:
        pass


# ------------------------------------------------------------------
# python-pptx helpers
# ------------------------------------------------------------------

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _delete_all_slides(prs: object) -> None:
    sld_id_lst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(f"{{{_NS_R}}}id")
        sld_id_lst.remove(sld_id)
        if r_id:
            try:
                prs.part.drop_rel(r_id)  # type: ignore[attr-defined]
            except Exception:
                pass


def _clone_slide(
    src_slide: object,
    src_prs: object,
    dst_prs: object,
) -> object:
    """Deep-clone a slide into dst_prs."""
    dst_layout = _match_layout(src_slide, dst_prs)
    new_slide = dst_prs.slides.add_slide(dst_layout)  # type: ignore[attr-defined]

    # Clear default placeholders
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Find spTree elements
    src_tree = src_slide._element.find(  # type: ignore[attr-defined]
        f"{{{_NS_P}}}cSld/{{{_NS_P}}}spTree"
    )
    dst_tree = new_slide._element.find(f"{{{_NS_P}}}cSld/{{{_NS_P}}}spTree")

    if src_tree is not None and dst_tree is not None:
        for child in list(dst_tree):
            tag = child.tag
            if not (tag.endswith("}nvGrpSpPr") or tag.endswith("}grpSpPr")):
                dst_tree.remove(child)

        for child in src_tree:
            tag = child.tag
            if tag.endswith("}nvGrpSpPr") or tag.endswith("}grpSpPr"):
                continue
            dst_tree.append(copy_module.deepcopy(child))

        rid_map = _copy_relationships(src_slide, new_slide)
        if rid_map:
            _remap_rids(dst_tree, rid_map)

    # Copy timing (animations)
    for tag_suffix in ("timing", "transition"):
        src_el = src_slide._element.find(  # type: ignore[attr-defined]
            f"{{{_NS_P}}}{tag_suffix}"
        )
        if src_el is not None:
            old = new_slide._element.find(f"{{{_NS_P}}}{tag_suffix}")
            if old is not None:
                new_slide._element.remove(old)
            new_slide._element.append(copy_module.deepcopy(src_el))

    # Copy notes
    try:
        if src_slide.has_notes_slide:  # type: ignore[attr-defined]
            txt = src_slide.notes_slide.notes_text_frame.text
            if txt.strip():
                new_slide.notes_slide.notes_text_frame.text = txt
    except Exception:
        pass

    return new_slide


def _match_layout(src_slide: object, dst_prs: object) -> object:
    try:
        name = src_slide.slide_layout.name  # type: ignore[attr-defined]
        for layout in dst_prs.slide_layouts:  # type: ignore[attr-defined]
            if layout.name == name:
                return layout
    except Exception:
        pass

    for fallback in ("Blank", "Leer", "Title Only"):
        for layout in dst_prs.slide_layouts:  # type: ignore[attr-defined]
            if layout.name == fallback:
                return layout

    return dst_prs.slide_layouts[0]  # type: ignore[attr-defined]


def _copy_relationships(src_slide: object, dst_slide: object) -> dict[str, str]:
    rid_map: dict[str, str] = {}
    skip = {
        f"{_NS_R.replace('/relationships', '')}/relationships/slideLayout",
        f"{_NS_R.replace('/relationships', '')}/relationships/notesSlide",
    }

    for rel in src_slide.part.rels.values():  # type: ignore[attr-defined]
        if rel.reltype in skip:
            continue
        try:
            new_rid = dst_slide.part.relate_to(  # type: ignore[attr-defined]
                rel.target_part, rel.reltype
            )
            if rel.rId != new_rid:
                rid_map[rel.rId] = new_rid
        except Exception:
            try:
                if rel.is_external:
                    new_rid = dst_slide.part.rels.get_or_add_ext_rel(  # type: ignore[attr-defined]
                        rel.reltype, rel.target_ref
                    )
                    if rel.rId != new_rid:
                        rid_map[rel.rId] = new_rid
            except Exception:
                pass

    return rid_map


def _remap_rids(element: object, rid_map: dict[str, str]) -> None:
    for attr_name in list(element.attrib.keys()):  # type: ignore[attr-defined]
        if element.attrib[attr_name] in rid_map:  # type: ignore[attr-defined]
            element.attrib[attr_name] = rid_map[  # type: ignore[attr-defined]
                element.attrib[attr_name]  # type: ignore[attr-defined]
            ]

    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for attr in ("embed", "link", "id"):
        full = f"{{{r_ns}}}{attr}"
        if full in element.attrib:  # type: ignore[attr-defined]
            old = element.attrib[full]  # type: ignore[attr-defined]
            if old in rid_map:
                element.attrib[full] = rid_map[old]  # type: ignore[attr-defined]

    for child in element:  # type: ignore[attr-defined]
        _remap_rids(child, rid_map)


# ------------------------------------------------------------------
# Public API
# ------------------------------------------------------------------


def merge_slides(
    config: MergeConfig,
    output_path: Path | str,
    *,
    backend: str = "auto",
) -> MergeResult:
    """Merge slides from multiple PPTX sources.

    Args:
        config: Merge configuration.
        output_path: Where to write the merged PPTX.
        backend: "auto", "com", or "pptx".

    Returns:
        MergeResult with output path, count, backend used.

    Raises:
        SlideMergeError: if merge fails.
    """
    output_path = Path(output_path)
    be = select_merge_backend(backend)

    # Validate sources exist
    for name, path in config.sources.items():
        if not Path(path).exists():
            raise SlideMergeError(f"Source '{name}' not found: {path}")
    if not Path(config.base).exists():
        raise SlideMergeError(f"Base template not found: {config.base}")

    try:
        if be == "com":
            impl = _COMBackend()
        else:
            impl = _PptxBackend()

        count = impl.merge(config, output_path)
        impl.close()
        logger.info(
            "Merged %d slides -> %s [%s]",
            count,
            output_path,
            be,
        )
        return MergeResult(
            output_path=output_path,
            slide_count=count,
            backend=be,
        )
    except Exception as exc:
        raise SlideMergeError(str(exc)) from exc
