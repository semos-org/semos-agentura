"""Convert Marp markdown to editable PPTX via pandoc.

Pipeline: Marp MD -> pandoc-flavored MD -> PPTX (via pandoc)
Optionally applies a corporate design template via COM automation.
"""

from __future__ import annotations

import logging
import re
import subprocess
from pathlib import Path

from ..exceptions import CompositionError, TemplateError

logger = logging.getLogger(__name__)

_DATA_DIR = Path(__file__).parent / "_data"


# ------------------------------------------------------------------
# Marp parser
# ------------------------------------------------------------------


def parse_marp(md_path: Path) -> list[dict]:
    """Parse Marp markdown into a list of slide dicts."""
    content = Path(md_path).read_text(encoding="utf-8")
    chunks = re.split(r"\n---\n", content)

    # Skip YAML frontmatter (first chunk starts with ---)
    if chunks and chunks[0].strip().startswith("---"):
        chunks = chunks[1:]

    return [_parse_slide(chunk) for chunk in chunks]


def _convert_html_columns(text: str) -> str:
    """Convert Marp HTML column divs to pandoc fenced div syntax.

    Input:  <div class="columns"><div>col1</div><div>col2</div></div>
    Output: :::::: {.columns}
            ::: {.column width="50%"}
            col1
            :::
            ::: {.column width="50%"}
            col2
            :::
            ::::::
    """
    # Match the outer columns div
    pat = re.compile(
        r'<div\s+class="columns">\s*'
        r"<div>(.*?)</div>\s*"
        r"<div>(.*?)</div>\s*"
        r"</div>",
        re.DOTALL,
    )

    def _replace(m: re.Match) -> str:
        col1 = m.group(1).strip()
        col2 = m.group(2).strip()
        return (
            ":::::::::::::: {.columns}\n"
            '::: {.column width="50%"}\n'
            f"{col1}\n"
            ":::\n"
            '::: {.column width="50%"}\n'
            f"{col2}\n"
            ":::\n"
            "::::::::::::::"
        )

    return pat.sub(_replace, text)


def _parse_slide(text: str) -> dict:
    """Parse a single slide chunk into structured data."""
    slide: dict = {
        "class": None,
        "title": None,
        "body_lines": [],
        "bg_images": [],
        "inline_images": [],
        "notes": None,
    }

    # Extract directives from comments
    directives: dict[str, str] = {}

    def _extract_directive(m: re.Match) -> str:
        directives[m.group(1).strip()] = m.group(2).strip()
        return ""

    text = re.sub(
        r"<!--\s*(_class|_paginate|_footer)\s*:\s*(.*?)\s*-->",
        _extract_directive,
        text,
        flags=re.DOTALL,
    )
    slide["class"] = directives.get("_class")

    # Extract speaker notes (remaining HTML comments)
    notes_parts: list[str] = []

    def _extract_notes(m: re.Match) -> str:
        lines = [ln.strip() for ln in m.group(1).strip().split("\n") if ln.strip()]
        if lines:
            notes_parts.append("\n".join(lines))
        return ""

    text = re.sub(r"<!--\s*(.*?)\s*-->", _extract_notes, text, flags=re.DOTALL)
    if notes_parts:
        slide["notes"] = "\n\n".join(notes_parts)

    # Background images: ![bg right:48% fit](path)
    bg_pat = re.compile(
        r"!\[bg\s*(left|right)?:?(\d+)?%?\s*"
        r"(?:fit|contain|cover)?\s*\]\(([^)]+)\)"
    )
    for m in bg_pat.finditer(text):
        slide["bg_images"].append(
            {
                "side": m.group(1) or "right",
                "pct": int(m.group(2)) if m.group(2) else 50,
                "path": m.group(3),
            }
        )
    text = bg_pat.sub("", text)

    # Inline images: ![w:880](path), ![w:800 center](path), ![](path)
    # Match any Marp directives in alt text (w:NNN, center, etc.)
    inline_pat = re.compile(r"!\[(?:[whWH]:\d+\s*)?(?:center\s*)?\]\(([^)]+)\)")
    for m in inline_pat.finditer(text):
        slide["inline_images"].append(m.group(1))
    text = inline_pat.sub("", text)

    # Convert Marp HTML columns to pandoc fenced div columns
    # <div class="columns"><div>col1</div><div>col2</div></div>
    text = _convert_html_columns(text)

    # Strip remaining HTML tags (Marp-specific: <br>, <div>, etc.)
    text = re.sub(r"</?(?:br|div|span)[^>]*>", "", text)

    # Parse title and body
    lines = text.strip().split("\n")
    body: list[str] = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("# ") and slide["title"] is None:
            slide["title"] = stripped[2:].strip()
        elif stripped.startswith("## ") and slide["title"] is None:
            slide["title"] = stripped[3:].strip()
        elif stripped:
            body.append(line)
        elif body:
            body.append("")

    while body and not body[-1].strip():
        body.pop()
    slide["body_lines"] = body
    return slide


# ------------------------------------------------------------------
# Pandoc markdown generation
# ------------------------------------------------------------------


def slide_to_pandoc(slide: dict) -> str:
    """Convert a parsed slide dict to pandoc markdown."""
    parts: list[str] = []

    # Title slide
    if slide["class"] == "title":
        title_lines: list[str] = []
        body_lines: list[str] = []
        if slide["title"]:
            title_lines.append(slide["title"])
        for line in slide["body_lines"]:
            if line.startswith("# ") and not body_lines:
                title_lines.append(line[2:])
            else:
                body_lines.append(line)
        if title_lines:
            parts.append(f"# {' - '.join(title_lines)}")
        for line in body_lines:
            parts.append(line)
        _add_notes(parts, slide)
        return "\n".join(parts)

    # Slide with background image -> two-column layout
    if slide["bg_images"]:
        if slide["title"]:
            parts.append(f"# {slide['title']}")
            parts.append("")

        imgs = slide["bg_images"]
        if len(imgs) == 2:
            parts.append(":::::::::::::: {.columns}")
            parts.append('::: {.column width="50%"}')
            parts.append(f"![]({imgs[0]['path']})")
            parts.append(":::")
            parts.append('::: {.column width="50%"}')
            parts.append(f"![]({imgs[1]['path']})")
            parts.append(":::")
            parts.append("::::::::::::::")
        elif len(imgs) == 1:
            img = imgs[0]
            text_pct = 100 - img["pct"]
            img_pct = img["pct"]
            parts.append(":::::::::::::: {.columns}")
            if img["side"] == "left":
                parts.append(f'::: {{.column width="{img_pct}%"}}')
                parts.append(f"![]({img['path']})")
                parts.append(":::")
                parts.append(f'::: {{.column width="{text_pct}%"}}')
                parts.extend(slide["body_lines"])
                parts.append(":::")
            else:
                parts.append(f'::: {{.column width="{text_pct}%"}}')
                parts.extend(slide["body_lines"])
                parts.append(":::")
                parts.append(f'::: {{.column width="{img_pct}%"}}')
                parts.append(f"![]({img['path']})")
                parts.append(":::")
            parts.append("::::::::::::::")
        _add_notes(parts, slide)
        return "\n".join(parts)

    # Regular slide
    if slide["title"]:
        parts.append(f"# {slide['title']}")
        parts.append("")

    has_body = any(ln.strip() for ln in slide["body_lines"])
    already_columns = any("{.columns}" in ln for ln in slide["body_lines"])
    has_table = any(ln.strip().startswith("|") and "---" in ln for ln in slide["body_lines"])

    # If body already has pandoc columns (from HTML div conversion),
    # just output as-is - don't double-wrap
    if already_columns:
        parts.extend(slide["body_lines"])
    elif slide["inline_images"] and has_body:
        # Image + text: wrap in columns to keep on same slide.
        # _polish_pptx restacks to rows if image is wide.
        parts.append(":::::::::::::: {.columns}")
        parts.append('::: {.column width="55%"}')
        for img_path in slide["inline_images"]:
            parts.append(f"![]({img_path})")
        parts.append(":::")
        parts.append('::: {.column width="45%"}')
        parts.extend(slide["body_lines"])
        parts.append(":::")
        parts.append("::::::::::::::")
    elif slide["inline_images"]:
        for img_path in slide["inline_images"]:
            parts.append(f"![]({img_path})")
            parts.append("")
    elif has_table and has_body:
        # Table + text: wrap in columns to force "Two Content"
        # layout instead of "Content with Caption".
        # Tables stay side-by-side (no restacking).
        table_lines: list[str] = []
        text_lines: list[str] = []
        in_table = False
        for ln in slide["body_lines"]:
            if ln.strip().startswith("|"):
                in_table = True
                table_lines.append(ln)
            else:
                if in_table and not ln.strip():
                    in_table = False
                text_lines.append(ln)
        if table_lines and any(ln.strip() for ln in text_lines):
            parts.append(":::::::::::::: {.columns}")
            parts.append('::: {.column width="60%"}')
            parts.extend(table_lines)
            parts.append(":::")
            parts.append('::: {.column width="40%"}')
            parts.extend(text_lines)
            parts.append(":::")
            parts.append("::::::::::::::")
        else:
            parts.extend(slide["body_lines"])
    else:
        parts.extend(slide["body_lines"])

    _add_notes(parts, slide)
    return "\n".join(parts)


def _get_image_aspect(img_path: str, source_dir: Path | str | None = None) -> float:
    """Get image width/height ratio. Returns 1.0 if unreadable."""
    candidates = [Path(img_path)]
    if source_dir:
        candidates.append(Path(source_dir) / img_path)
        candidates.append(Path(source_dir) / Path(img_path).name)

    for p in candidates:
        if not p.exists():
            continue
        try:
            # Read PNG/JPEG dimensions without Pillow
            with open(p, "rb") as f:
                header = f.read(32)
            if header[:8] == b"\x89PNG\r\n\x1a\n":
                # PNG: width at bytes 16-20, height at 20-24
                import struct

                w, h = struct.unpack(">II", header[16:24])
                return w / h if h else 1.0
            if header[:2] in (b"\xff\xd8",):
                # JPEG: need to parse markers, fallback to Pillow
                try:
                    from PIL import Image

                    with Image.open(p) as im:
                        w, h = im.size
                        return w / h if h else 1.0
                except ImportError:
                    return 1.0
            if header[:4] == b"GIF8":
                import struct

                w, h = struct.unpack("<HH", header[6:10])
                return w / h if h else 1.0
        except Exception:
            pass
    return 1.0  # unknown -> assume square


def _set_pos(shape: object, left: int, top: int, width: int, height: int) -> None:
    """Set all 4 positional properties at once.

    python-pptx resets inherited values to 0 when you set a single
    property on a placeholder without an explicit xfrm. This helper
    creates the xfrm element by writing all 4 properties, ensuring
    none are lost.
    """
    from lxml import etree

    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    sp = shape._element

    # Find xfrm: in spPr for shapes, directly under element for
    # graphicFrame (tables, charts)
    xfrm = None
    for parent_tag in (f"{{{ns_p}}}spPr", f"{{{ns_a}}}spPr"):
        parent = sp.find(parent_tag)
        if parent is not None:
            xfrm = parent.find(f"{{{ns_a}}}xfrm")
            if xfrm is None:
                xfrm = etree.SubElement(parent, f"{{{ns_a}}}xfrm")
            break
    if xfrm is None:
        # graphicFrame: xfrm is p:xfrm directly under element
        xfrm = sp.find(f"{{{ns_p}}}xfrm")
    if xfrm is None:
        # Last resort: fall back to python-pptx
        shape.left = int(left)
        shape.top = int(top)
        shape.width = int(width)
        shape.height = int(height)
        return

    off = xfrm.find(f"{{{ns_a}}}off")
    ext = xfrm.find(f"{{{ns_a}}}ext")
    if off is None:
        off = etree.SubElement(xfrm, f"{{{ns_a}}}off")
    if ext is None:
        ext = etree.SubElement(xfrm, f"{{{ns_a}}}ext")
    off.set("x", str(int(left)))
    off.set("y", str(int(top)))
    ext.set("cx", str(int(width)))
    ext.set("cy", str(int(height)))


def _get_pptx_image_aspect(placeholder: object) -> float:
    """Get aspect ratio of image inside a PPTX placeholder.

    Reads the picture's embedded extent or the image blob dimensions.
    Returns width/height ratio, defaults to 1.0 if unreadable.
    """
    try:
        # Look for a pic element with blipFill
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        for pic in placeholder._element.iter(f"{{{ns_a}}}blip"):
            # Get the relationship ID to find the image part
            embed = pic.get(f"{{{ns_r}}}embed")
            if embed:
                try:
                    img_part = placeholder.part.related_parts[embed]
                    blob = img_part.blob
                    # PNG header
                    if blob[:8] == b"\x89PNG\r\n\x1a\n":
                        import struct

                        w, h = struct.unpack(">II", blob[16:24])
                        if h:
                            return w / h
                except Exception:
                    pass

        # Fallback: check the extent (cx, cy) of the pic element
        for ext in placeholder._element.iter(f"{{{ns_a}}}ext"):
            cx = ext.get("cx")
            cy = ext.get("cy")
            if cx and cy:
                cx_val, cy_val = int(cx), int(cy)
                if cy_val:
                    return cx_val / cy_val
    except Exception:
        pass
    return 1.0


def _add_notes(parts: list[str], slide: dict) -> None:
    if slide["notes"]:
        parts.append("")
        parts.append("::: notes")
        parts.append(slide["notes"])
        parts.append(":::")


# ------------------------------------------------------------------
# PPTX post-processing
# ------------------------------------------------------------------


def _polish_pptx(pptx_path: Path, slides_data: list[dict]) -> None:
    """Widen titles and fix zero-height placeholders.

    Conservative: only touches what we can safely identify.
    Leaves image+text slides in pandoc's default layout.

    IMPORTANT: when setting any positional property on a placeholder
    that inherits from the slide layout, python-pptx resets the others
    to 0. Always save and restore all 4 properties together.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        logger.debug("python-pptx not installed, skipping polish")
        return

    prs = Presentation(str(pptx_path))
    sw = prs.slide_width
    pad = Inches(0.5)
    title_bottom = Inches(1.2)
    max_content_h = Inches(5.8)

    # Pandoc may produce more slides than we parsed (table/column
    # splits). Pad slides_data so every slide gets processed.
    padded = list(slides_data) + [{"class": None, "bg_images": [], "inline_images": []}] * max(
        0, len(prs.slides) - len(slides_data)
    )

    for slide, data in zip(prs.slides, padded, strict=False):
        title_sh = None
        content_phs: list = []

        for sh in slide.shapes:
            if not sh.is_placeholder:
                continue
            idx = sh.placeholder_format.idx
            if idx == 0:
                title_sh = sh
            elif idx in (1, 2):
                content_phs.append(sh)

        if title_sh:
            t, h = title_sh.top, title_sh.height
            _set_pos(
                title_sh,
                pad,
                t if t else Inches(0.2),
                sw - 2 * pad,
                h if h and h >= Inches(0.5) else Inches(1.0),
            )
            # Normalize title font across all layouts
            if title_sh.has_text_frame:
                for para in title_sh.text_frame.paragraphs:
                    for run in para.runs:
                        from pptx.util import Pt

                        if data["class"] == "title":
                            run.font.size = Pt(36)
                        else:
                            run.font.size = Pt(28)
                        run.font.bold = False
            if data["class"] == "title":
                title_sh.top = Inches(2.0)
                title_sh.height = Inches(1.5)

        if data["class"] == "title":
            for cs in content_phs:
                cs.left = Inches(2.0)
                cs.width = sw - Inches(4.0)
                cs.top = Inches(3.8)
                cs.height = Inches(2.5)
            continue

        # --- Step 1: Detect media+text and try restacking to rows ---
        restacked = False
        free_pics = [
            sh
            for sh in slide.shapes
            if not sh.is_placeholder and any(c.tag.endswith("}pic") for c in sh._element.iter())
        ]

        media_sh = txt_sh = None
        for cs in content_phs:
            has_media = any(
                c.tag.endswith("}pic") or c.tag.endswith("}tbl") or c.tag.endswith("}graphic")
                for c in cs._element.iter()
            )
            if has_media and not media_sh:
                media_sh = cs
            elif not txt_sh:
                txt_sh = cs
        if not media_sh and free_pics:
            media_sh = free_pics[0]
        if not txt_sh and len(content_phs) == 1 and media_sh:
            txt_sh = content_phs[0]
            if txt_sh is media_sh:
                txt_sh = None

        if media_sh and txt_sh:
            # Get aspect ratio: images from blob, tables from
            # placeholder dimensions
            aspect = _get_pptx_image_aspect(media_sh)
            if aspect == 1.0 and media_sh.width and media_sh.height:
                # Fallback: use shape dimensions (works for tables)
                aspect = media_sh.width / media_sh.height
            if aspect >= 1.3:
                avail_w = sw - 2 * pad
                img_h = int(avail_w / aspect)
                max_img_h = int(max_content_h * 0.65)
                if img_h > max_img_h:
                    img_h = max_img_h
                gap = Inches(0.1)
                txt_h = int(max_content_h) - img_h - gap
                if txt_h < Inches(1.2):
                    txt_h = Inches(1.2)
                img_w = int(min(avail_w, img_h * aspect))
                img_left = pad + (avail_w - img_w) // 2

                _set_pos(media_sh, img_left, title_bottom, img_w, img_h)
                _set_pos(txt_sh, pad, int(title_bottom) + img_h + gap, sw - 2 * pad, txt_h)
                restacked = True

        # --- Step 2: Fix remaining placeholders (not restacked) ---
        if not restacked:
            for cs in content_phs:
                t, h, left, w = cs.top, cs.height, cs.left, cs.width
                needs_fix = (not t or t < title_bottom) or (not h or h < Inches(0.5)) or (h > max_content_h)
                if needs_fix:
                    cs.left = left if left else pad
                    cs.top = title_bottom
                    cs.width = w if w else sw - 2 * pad
                    cs.height = h if h and Inches(0.5) <= h <= max_content_h else max_content_h

        # Text-only slides: widen to full width
        has_images = data["bg_images"] or data["inline_images"]
        if not has_images:
            for cs in content_phs:
                t, h = cs.top, cs.height
                cs.left = pad
                cs.top = t
                cs.width = sw - 2 * pad
                cs.height = h

    prs.save(str(pptx_path))
    logger.info("Polished PPTX layout: %s", pptx_path)


# ------------------------------------------------------------------
# Template application
# ------------------------------------------------------------------


def _apply_template_com(pptx_path: Path, template_path: Path) -> None:
    """Apply template via PowerPoint COM automation (Windows)."""
    import win32com.client

    ppt = win32com.client.Dispatch("PowerPoint.Application")
    pres = ppt.Presentations.Open(str(pptx_path.resolve()))
    pres.ApplyTemplate(str(template_path.resolve()))
    pres.Save()
    pres.Close()
    logger.info("Applied template via COM: %s", template_path)


def _apply_template_uno(pptx_path: Path, template_path: Path) -> None:
    """Apply template via LibreOffice UNO (EXPERIMENTAL stub).

    Transfers placeholder structure only, not visual branding.
    LO's OOXML importer doesn't expose master branding shapes.
    """
    raise NotImplementedError(
        "UNO template application is experimental and transfers "
        "placeholder structure only, not visual branding. "
        "Use COM backend (PowerPoint) for full branding."
    )


def _apply_template_docker(pptx_path: Path, template_path: Path) -> None:
    """Apply template via Docker+LO UNO (EXPERIMENTAL stub).

    Same limitations as UNO: no visual branding transfer.
    """
    raise NotImplementedError(
        "Docker template application is experimental. Use COM backend (PowerPoint) for full branding."
    )


def apply_template(
    pptx_path: Path,
    template_path: Path,
    backend: str = "auto",
) -> None:
    """Apply a design template to an existing PPTX.

    Args:
        pptx_path: Path to the PPTX to brand.
        template_path: Path to the template PPTX.
        backend: "auto", "com", "uno", or "docker".

    Raises:
        TemplateError: if no backend succeeds.
    """
    pptx_path = Path(pptx_path)
    template_path = Path(template_path)

    if backend == "auto":
        backends = ["com", "uno", "docker"]
    else:
        backends = [backend]

    for be in backends:
        try:
            if be == "com":
                _apply_template_com(pptx_path, template_path)
                return
            elif be == "uno":
                _apply_template_uno(pptx_path, template_path)
                return
            elif be == "docker":
                _apply_template_docker(pptx_path, template_path)
                return
        except ImportError:
            continue
        except NotImplementedError as exc:
            logger.debug("%s backend: %s", be, exc)
            continue
        except Exception as exc:
            logger.warning("%s backend failed: %s", be, exc)
            continue

    raise TemplateError("No template backend available. PowerPoint COM (Windows) is required for full branding.")


# ------------------------------------------------------------------
# Main pipeline
# ------------------------------------------------------------------


def default_reference() -> Path:
    """Return the bundled widescreen pandoc reference PPTX."""
    ref = _DATA_DIR / "pandoc_reference.pptx"
    if not ref.exists():
        raise CompositionError(f"Bundled pandoc reference not found: {ref}")
    return ref


def _is_marp(md_path: Path) -> bool:
    """Detect whether a markdown file is Marp (has marp: true).

    Checks for ``marp: true`` in the YAML frontmatter (between
    opening and closing ``---`` lines).
    """
    try:
        text = md_path.read_text(encoding="utf-8")
        if not text.strip().startswith("---"):
            return False
        # Find closing --- (skip the opening one)
        end = text.find("\n---", 4)
        if end < 0:
            return False
        frontmatter = text[3:end]
        return bool(
            re.search(
                r"^\s*marp\s*:\s*true",
                frontmatter,
                re.MULTILINE,
            )
        )
    except Exception:
        return False


def compose_editable_slides(
    md_path: Path | str,
    output_path: Path | str,
    *,
    pandoc_path: Path | str = "pandoc",
    reference_doc: Path | str | None = None,
    template: Path | str | None = None,
    template_backend: str = "auto",
) -> Path:
    """Convert markdown slides to editable PPTX.

    Auto-detects the input format:
    - Marp markdown (has ``marp: true`` in frontmatter):
      parsed and converted to pandoc markdown, then to PPTX.
    - Plain pandoc markdown: passed directly to pandoc.

    Args:
        md_path: Path to markdown slide file (Marp or pandoc).
        output_path: Where to write the output PPTX.
        pandoc_path: Path to pandoc executable.
        reference_doc: Pandoc reference PPTX for theme/layout.
            Defaults to bundled widescreen reference.
        template: Optional design template PPTX to apply after
            generation (requires PowerPoint COM or similar).
        template_backend: Backend for template application.

    Returns:
        Path to the output PPTX.

    Raises:
        CompositionError: if pandoc fails or polish crashes.
        TemplateError: if template application fails.
    """
    md_path = Path(md_path)
    output_path = Path(output_path)

    # Resolve reference doc
    if reference_doc is None:
        reference_doc = default_reference()
    else:
        reference_doc = Path(reference_doc)

    if _is_marp(md_path):
        slides = parse_marp(md_path)
        logger.info(
            "Marp detected: parsed %d slides from %s",
            len(slides),
            md_path,
        )
        source_dir = md_path.parent
        for slide in slides:
            slide["_source_dir"] = source_dir

        # Build pandoc markdown from parsed Marp
        pandoc_parts: list[str] = []
        for i, slide in enumerate(slides):
            if i > 0:
                pandoc_parts.append("")
            pandoc_parts.append(slide_to_pandoc(slide))
        pandoc_md = "\n".join(pandoc_parts) + "\n"

        pandoc_input = output_path.with_suffix(".pandoc.md")
        pandoc_input.write_text(pandoc_md, encoding="utf-8")
    else:
        logger.info("Pandoc markdown detected: %s", md_path)
        slides = []
        pandoc_input = md_path

    # Run pandoc - choose output format from file extension
    out_ext = output_path.suffix.lower()
    if out_ext == ".pdf":
        pandoc_format = "beamer"
    elif out_ext == ".html":
        pandoc_format = "revealjs"
    else:
        pandoc_format = "pptx"

    cmd = [
        str(pandoc_path),
        str(pandoc_input),
        "-t",
        pandoc_format,
        f"--resource-path={md_path.parent}",
        "-o",
        str(output_path),
    ]
    if pandoc_format == "pptx":
        cmd.append(f"--reference-doc={reference_doc}")
    logger.info("Running: %s", " ".join(cmd))
    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    if result.returncode != 0:
        raise CompositionError(f"pandoc failed: {result.stderr.strip()}")

    # Post-process layout (Marp slides only - pandoc md is
    # already in pandoc's native format, polish would misalign)
    if slides:
        try:
            _polish_pptx(output_path, slides)
        except Exception as exc:
            logger.warning("Polish step failed: %s", exc)

    logger.info("Saved editable PPTX: %s", output_path)

    # Apply template if requested
    if template:
        apply_template(output_path, Path(template), backend=template_backend)

    return output_path
